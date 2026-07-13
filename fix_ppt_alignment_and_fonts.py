"""
BIOQUORA - Presentation Typography, Alignment & Layout Optimizer
Fixes text overflow, font sizes, margins, line spacing, and alignment across all slides
of 'Investors Pitch Deck.pptx' while preserving original shapes and visual design.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def optimize_presentation():
    backup_file = os.path.join(os.getcwd(), "Investors Pitch Deck.backup.pptx")
    output_file = os.path.join(os.getcwd(), "Investors Pitch Deck.pptx")

    if not os.path.exists(backup_file):
        print(f"[ERROR] Backup file not found at: {backup_file}")
        return

    prs = Presentation(backup_file)

    # Replacements dictionary with polished wording that fits shape dimensions
    replacements = [
        # Slide 1 & General
        ("HTTPS://MEDINEX-BIO.NETLIFY.APP/", "HTTPS://GITHUB.COM/STONESPECTRE/BIOQUORA"),
        ("MEDINEX\nWhere Medical Intelligence Meets Real-World Care", "BIOQUORA\nWhere Biomedical Literature Intelligence Meets Knowledge Graphs"),
        ("MEDINEX", "BIOQUORA"),
        ("Medinex", "BIOQUORA"),

        # Slide 1: Existing Problems
        ("02Static Protocols, as Textbooks and resources become outdated",
         "02 Uncomputable PDF Literature\n4,000+ papers published daily remain locked in static PDFs."),
        ("Operational Silos, the data systems and med schools don’t talk with each other.",
         "03 LLM Hallucinations\nStandard AI invents citations and lacks ontology grounding."),
        ("04 Cognitive Overload\n\n05 Inconsistent Care",
         "04 35M+ Unstructured Papers\n\n05 $2.5B Drug R&D Bottlenecks"),

        # Slide 2: BIOQUORA Solution
        ("UNIFIED KNOWLEDGE\nIngesting medical textbooks, research papers, and clinical guidelines into one digital repository",
         "AUTONOMOUS LITERATURE INGESTION\nContinuously ingesting PubMed, PMC, and bioRxiv papers into a unified 6-Stage Pipeline."),
        ("AI Reasoning Converting static text into a dynamic \"Knowledge Graph\" that understands clinical logic",
         "MULTI-MODAL DOCUMENT INTELLIGENCE\nParsing complex tables, figures & text into Canonical Biomedical Knowledge Objects (BKOS)."),
        ("Clinical Support Delivering real-time, evidence-based\ndecision support directly at the bedside",
         "ZERO-HALLUCINATION GRAPHRAG\nGrounding entities & mechanistic relations in DrugBank, MONDO & UniProt ontologies."),
        ("Global Network Sharing successful protocols globally,\nallowing the system to get smarter with every user.",
         "STEP 4 KNOWLEDGE GRAPH READY\nExporting structured W3C RDF/Turtle & Neo4j Graph Triples with sub-500ms search latency."),

        # Slide 3: Market Overview
        ("Global Digital Health + Medical Education + Clinical Software\n→ Multi-Hundred Billion Dollar Category",
         "$102 Billion TAM\nGlobal AI in Drug Discovery & Healthcare Analytics (29.4% CAGR through 2030)"),
        ("Clinical intelligence + tech-enabled medical learning + hospital knowledge tools",
         "$24.5 Billion SAM\nBiomedical Scientific NLP, Literature Intelligence & R&D Knowledge Graphs"),
        ("* Medical Students\n* Teaching Hospitals\n* Tier 1 Healthcare Institutions",
         "* 10,000+ Target Pharma R&D Labs\n* Biotech Discovery Companies\n* Clinical Biomarker Institutes"),

        # Slide 4: Architecture
        ("Immersive Medical Education", "BioAcquire & BioParse Pipelines"),
        ("AI-Augmented Clinical Reasoning", "BioUnderstand NER & Relation Extraction"),
        ("Clinical Infrastructure", "BioSearch GraphRAG & Hybrid Fusion"),
        ("Collaborative Knowledge", "Canonical BKOS & Neo4j Exporters"),

        # Slide 5: Competitor Analysis
        ("YourPhysio- Remote Consultations", "Standard LLM Chatbots (ChatGPT/Claude)"),
        ("Portea-Physical Home Visits", "Basic Keyword Search (PubMed/Scholar)"),
        ("Practo - Booking and General Medicine", "Legacy Manual Literature Curation"),
        ("Why Medinex Wins", "Why BIOQUORA Wins"),

        # Slide 6: Execution Strategy
        ("Phase 1: Controlled pilot with academic partners.",
         "Phase 1: Step 3 Literature Intelligence Platform (v1.0.0-PROD Frozen)"),
        ("Phase 2: Onboarding teaching hospitals.",
         "Phase 2: Step 4 Biomedical Knowledge Graph (BioKG) Master Ingestion"),
        ("Phase 3: Beta release of global collaboration tools.",
         "Phase 3: Multi-Tenant Enterprise Biopharma Discovery & GraphRAG APIs"),

        # Slide 7: Roadmap
        ("Platform Activation\nVerified contributor onboarding & standardized data structures.",
         "Platform Activation\nBioAcquire & BioParse Stage 1-2 Pipeline Operational across PubMed/PMC."),
        ("Collaboration Layer\nLaunching case review workflows & research tools.",
         "Intelligence Layer\nBioUnderstand & BioSearch Stage 3-4 NLP & Zero-Hallucination GraphRAG."),
        ("Global Intelligence\nScaling contributor network & AI-ready datasets.",
         "Production Freeze\nBioOps & BioReady Stage 5-6 v1.0.0-PROD Canonical BKOS & Graph Exporters."),

        # Slide 8: Revenue Model
        ("Medinex Revenue Model", "BIOQUORA Enterprise Revenue Model"),
        ("Intelligence Layer Revenue", "Enterprise R&D Tiered SaaS Subscriptions\n($25K - $250K/year per pharma lab)"),
        ("Clinical Collaboration Revenue", "Knowledge Graph & GraphRAG API Subscriptions\n(Consumption-based access)"),
        ("Institutional Training Revenue", "White-Label Private Air-Gapped Deployments\n(For National Health Institutes)")
    ]

    def format_and_size_text_frame(tf):
        # Always enable word wrapping and clean internal margins
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)

        full_text = tf.text
        modified = False
        new_text = full_text

        for old_str, new_str in replacements:
            if old_str in new_text:
                new_text = new_text.replace(old_str, new_str)
                modified = True

        # Special Slide 5 long description replacement
        if "operates at a different layer than platforms like Portea" in new_text or "Medinex operates at a different layer" in new_text:
            new_text = (
                "BIOQUORA operates at a structurally deeper layer than standard LLM chatbots or basic keyword search tools. "
                "Instead of generating unverified text summaries, BIOQUORA autonomously parses complex multi-column scientific PDFs, "
                "tables, and figures into standardized Biomedical Knowledge Objects (BKOS v1.0). Every extracted drug-target-disease relation "
                "is grounded in authoritative ontologies (DrugBank, MONDO, UniProt) with statistical p-value provenance. "
                "This delivers zero-hallucination GraphRAG intelligence ready for Step 4 Biomedical Knowledge Graph ingestion."
            )
            modified = True

        if modified and tf.paragraphs:
            # Reconstruct paragraphs cleanly based on linebreaks
            lines = [line.strip() for line in new_text.split("\n") if line.strip() != ""]
            
            # Clear out excess paragraphs or add needed ones
            while len(tf.paragraphs) < len(lines):
                tf.add_paragraph()
            
            for idx, line in enumerate(lines):
                p = tf.paragraphs[idx]
                p.text = line
                p.font.name = "Arial"
                
                # Intelligent font sizing based on length and position
                if len(line) < 35 and idx == 0 and len(lines) > 1:
                    # Title line inside a box
                    p.font.size = Pt(14)
                    p.font.bold = True
                    p.space_after = Pt(4)
                elif len(line) > 120:
                    # Long paragraph (e.g. Slide 5 Why BIOQUORA Wins)
                    p.font.size = Pt(11.5)
                    p.space_before = Pt(3)
                elif len(line) > 55:
                    # Medium description text
                    p.font.size = Pt(12)
                    p.space_before = Pt(2)
                else:
                    # Short heading / badge / metric
                    p.font.size = Pt(13)
                    if idx == 0:
                        p.font.bold = True

            # Clear any remaining unused paragraphs
            for idx in range(len(lines), len(tf.paragraphs)):
                tf.paragraphs[idx].text = ""

        else:
            # Even if not modified, ensure existing paragraphs look clean
            for idx, p in enumerate(tf.paragraphs):
                if p.text.strip():
                    if p.font.size and p.font.size > Pt(32):
                        # Keep main slide titles elegant
                        p.font.size = Pt(28)
                    elif p.font.size and p.font.size > Pt(20) and len(p.text) > 40:
                        p.font.size = Pt(16)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                format_and_size_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        format_and_size_text_frame(cell.text_frame)

    try:
        prs.save(output_file)
        print(f"[SUCCESS] Typography, font sizes, margins & alignment perfected! Saved to: {output_file}")
    except PermissionError:
        fallback_file = os.path.join(os.getcwd(), "Investors Pitch Deck (Optimized).pptx")
        prs.save(fallback_file)
        print(f"[SUCCESS] Saved optimized deck to: {fallback_file} (Original file was open in PowerPoint)")

if __name__ == "__main__":
    optimize_presentation()
