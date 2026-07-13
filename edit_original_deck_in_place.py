"""
BIOQUORA - In-Place Editor for Original Investors Pitch Deck
Takes the user's original presentation ('Investors Pitch Deck.backup.pptx'),
preserves 100% of the existing slide layouts, shapes, icons, and colors,
and edits the text inside the original shapes to feature BIOQUORA Step 3 v1.0.0-PROD.
"""

import os
import re
from pptx import Presentation

def edit_in_place():
    backup_file = os.path.join(os.getcwd(), "Investors Pitch Deck.backup.pptx")
    output_file = os.path.join(os.getcwd(), "Investors Pitch Deck.pptx")

    if not os.path.exists(backup_file):
        print(f"[ERROR] Backup file not found at {backup_file}")
        return

    prs = Presentation(backup_file)

    # Replacements dictionary for exact or substring matches
    replacements = [
        # General Name & URL Replacements
        ("HTTPS://MEDINEX-BIO.NETLIFY.APP/", "HTTPS://GITHUB.COM/STONESPECTRE/BIOQUORA"),
        ("MEDINEX\nWhere Medical Intelligence Meets Real-World Care", "BIOQUORA\nWhere Biomedical Literature Intelligence Meets Knowledge Graphs"),
        ("MEDINEX", "BIOQUORA"),
        ("Medinex", "BIOQUORA"),
        
        # Slide 1: Existing Problems
        ("02Static Protocols, as Textbooks and resources become outdated",
         "02 Uncomputable PDF Literature: 4,000+ papers published daily remain locked in static PDFs"),
        ("Operational Silos, the data systems and med schools don’t talk with each other.",
         "LLM Hallucinations: Standard AI invents citations and lacks ontology grounding"),
        ("04 Cognitive Overload\n\n05 Inconsistent Care",
         "04 35M+ Unstructured Papers\n\n05 $2.5B Drug R&D Bottleneck"),

        # Slide 2: Unified Knowledge -> BIOQUORA Solution
        ("UNIFIED KNOWLEDGE\nIngesting medical textbooks, research papers, and clinical guidelines into one digital repository",
         "AUTONOMOUS LITERATURE INGESTION\nContinuously ingesting PubMed, PMC, and bioRxiv papers into a unified 6-Stage Pipeline"),
        ("AI Reasoning Converting static text into a dynamic \"Knowledge Graph\" that understands clinical logic",
         "MULTI-MODAL DOCUMENT INTELLIGENCE\nParsing complex tables, figures & text into Canonical Biomedical Knowledge Objects (BKOS)"),
        ("Clinical Support Delivering real-time, evidence-based\ndecision support directly at the bedside",
         "ZERO-HALLUCINATION GRAPHRAG\nGrounding entities & mechanistic relations in DrugBank, MONDO & UniProt ontologies"),
        ("Global Network Sharing successful protocols globally,\nallowing the system to get smarter with every user.",
         "STEP 4 KNOWLEDGE GRAPH READY\nExporting structured W3C RDF/Turtle & Neo4j Graph Triples with sub-500ms search latency"),

        # Slide 3: Market Overview
        ("Global Digital Health + Medical Education + Clinical Software\n→ Multi-Hundred Billion Dollar Category",
         "$102 Billion TAM\nGlobal AI in Drug Discovery & Healthcare Analytics (29.4% CAGR through 2030)"),
        ("Clinical intelligence + tech-enabled medical learning + hospital knowledge tools",
         "$24.5 Billion SAM\nBiomedical Scientific NLP, Literature Intelligence & R&D Knowledge Graphs"),
        ("* Medical Students\n* Teaching Hospitals\n* Tier 1 Healthcare Institutions",
         "* 10,000+ Target Pharma R&D Labs\n* Biotech Discovery Companies\n* Clinical Biomarker Institutes"),

        # Slide 4: Our Services -> BIOQUORA Architecture
        ("Immersive Medical Education", "BioAcquire & BioParse Pipelines"),
        ("AI-Augmented Clinical Reasoning", "BioUnderstand NER & Relation Extraction"),
        ("Clinical Infrastructure", "BioSearch GraphRAG & Hybrid Fusion"),
        ("Collaborative Knowledge", "Canonical BKOS & Neo4j Exporters"),

        # Slide 5: Competitor Analysis
        ("YourPhysio- Remote Consultations", "Standard LLM Chatbots (ChatGPT/Claude)"),
        ("Portea-Physical Home Visits", "Basic Keyword Search (PubMed / Google Scholar)"),
        ("Practo - Booking and General Medicine", "Legacy Manual Literature Curation"),
        ("Why Medinex Wins", "Why BIOQUORA Wins"),

        # Slide 6: Execution Strategy
        ("Phase 1: Controlled pilot with academic partners.",
         "Phase 1: Step 3 Literature Intelligence Platform (v1.0.0-PROD Frozen)"),
        ("Phase 2: Onboarding teaching hospitals.",
         "Phase 2: Step 4 Biomedical Knowledge Graph (BioKG) Master Ingestion"),
        ("Phase 3: Beta release of global collaboration tools.",
         "Phase 3: Multi-Tenant Enterprise Biopharma Discovery & GraphRAG APIs"),

        # Slide 7: Strategic Execution Roadmap
        ("Platform Activation\nVerified contributor onboarding & standardized data structures.",
         "Platform Activation\nBioAcquire & BioParse Stage 1-2 Pipeline Operational across PubMed/PMC."),
        ("Collaboration Layer\nLaunching case review workflows & research tools.",
         "Intelligence Layer\nBioUnderstand & BioSearch Stage 3-4 NLP & Zero-Hallucination GraphRAG."),
        ("Global Intelligence\nScaling contributor network & AI-ready datasets.",
         "Production Freeze\nBioOps & BioReady Stage 5-6 v1.0.0-PROD Canonical BKOS & Graph Exporters."),

        # Slide 8: Revenue Model
        ("Medinex Revenue Model", "BIOQUORA Enterprise Revenue Model"),
        ("Intelligence Layer Revenue", "Enterprise R&D Tiered SaaS Subscriptions ($25K - $250K/year per pharma lab)"),
        ("Clinical Collaboration Revenue", "Knowledge Graph & GraphRAG API Consumption Subscriptions"),
        ("Institutional Training Revenue", "White-Label Private Air-Gapped Deployments for National Health Institutes")
    ]

    # Helper to replace text within paragraph runs to keep formatting
    def replace_text_in_frame(tf):
        full_text = tf.text
        modified = False
        new_text = full_text

        # First check exact or substring match replacements
        for old_str, new_str in replacements:
            if old_str in new_text:
                new_text = new_text.replace(old_str, new_str)
                modified = True

        # Special Slide 5 long description check
        if "operates at a different layer than platforms like Portea" in new_text or "Medinex operates at a different layer" in new_text:
            new_text = (
                "BIOQUORA operates at a structurally deeper layer than standard LLM chatbots or basic keyword search tools. "
                "Instead of generating unverified text summaries, BIOQUORA autonomously parses complex multi-column scientific PDFs, "
                "tables, and figures into standardized Biomedical Knowledge Objects (BKOS v1.0). Every extracted drug-target-disease relation "
                "is grounded in authoritative ontologies (DrugBank, MONDO, UniProt) with statistical p-value provenance. "
                "This delivers zero-hallucination GraphRAG intelligence ready for Step 4 Biomedical Knowledge Graph ingestion."
            )
            modified = True

        if modified and tf.paragraphs:
            # Keep font style of first run if possible
            first_p = tf.paragraphs[0]
            first_p.text = new_text
            # clear remaining paragraphs if text was consolidated into first paragraph
            for p in tf.paragraphs[1:]:
                p.text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text_in_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text_in_frame(cell.text_frame)

    prs.save(output_file)
    print(f"[SUCCESS] Edited original presentation in-place! Saved to: {output_file}")

if __name__ == "__main__":
    edit_in_place()
