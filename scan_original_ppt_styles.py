"""
BIOQUORA - Exact Presentation Style & Font Scanner
Scans 'Investors Pitch Deck.backup.pptx' slide by slide, shape by shape,
reporting exact font names, font sizes, colors, bold/italic styles, and alignments.
"""

import os
from pptx import Presentation

def scan_styles():
    backup_file = os.path.join(os.getcwd(), "Investors Pitch Deck.backup.pptx")
    if not os.path.exists(backup_file):
        print(f"[ERROR] Backup file not found: {backup_file}")
        return

    prs = Presentation(backup_file)
    print(f"=== SCANNING ORIGINAL PRESENTATION STYLES ({len(prs.slides)} SLIDES) ===\n")

    for s_idx, slide in enumerate(prs.slides):
        print(f"--- SLIDE {s_idx+1} ---")
        for sh_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                tf = shape.text_frame
                for p_idx, p in enumerate(tf.paragraphs):
                    text_snippet = p.text.strip()[:40].replace('\n', ' ')
                    if not text_snippet:
                        continue
                    font_name = p.font.name
                    font_size = p.font.size.pt if (p.font and p.font.size) else "Default"
                    font_bold = p.font.bold if p.font else "Default"
                    font_color = None
                    if p.font and p.font.color and p.font.color.type:
                        try:
                            font_color = str(p.font.color.rgb)
                        except Exception:
                            font_color = "Theme/Inherited"
                    # also inspect first run if paragraph font properties are None
                    if len(p.runs) > 0:
                        r = p.runs[0]
                        if not font_name and r.font: font_name = r.font.name
                        if font_size == "Default" and r.font and r.font.size: font_size = r.font.size.pt
                        if font_bold == "Default" and r.font: font_bold = r.font.bold
                        if not font_color and r.font and r.font.color:
                            try: font_color = str(r.font.color.rgb)
                            except Exception: pass
                    print(f"  Shape {sh_idx+1} | P {p_idx+1}: '{text_snippet}' -> Font: {font_name} | Size: {font_size} pt | Bold: {font_bold} | Color: {font_color}")

if __name__ == "__main__":
    scan_styles()
